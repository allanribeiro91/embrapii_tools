import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def pdf_to_ppt_image(pdf_file_path, output_ppt_path, dpi=600):
    # Criar uma pasta temporária para armazenar as imagens convertidas
    temp_folder = "temp_images"
    if not os.path.exists(temp_folder):
        os.makedirs(temp_folder)

    try:
        # Abrir o PDF com PyMuPDF
        print("Convertendo PDF em imagens...")
        pdf_document = fitz.open(pdf_file_path)

        # Criar uma apresentação PPT vazia com formato 16:9
        presentation = Presentation()
        presentation.slide_width = Inches(13.33)  # Largura para 16:9 (13.33 polegadas)
        presentation.slide_height = Inches(7.5)   # Altura para 16:9 (7.5 polegadas)

        # Converter cada página do PDF em uma imagem e adicionar ao PPT
        for page_number in range(len(pdf_document)):
            print(f"Adicionando página {page_number + 1} ao PPT...")
            page = pdf_document.load_page(page_number)

            # Renderizar a página como imagem
            pix = page.get_pixmap(dpi=dpi)
            image_path = os.path.join(temp_folder, f"page_{page_number + 1}.png")
            pix.save(image_path)

            # Adicionar a imagem ao slide
            slide_layout = presentation.slide_layouts[6]  # Layout vazio
            slide = presentation.slides.add_slide(slide_layout)

            left = top = Inches(0)
            slide_width = presentation.slide_width
            slide_height = presentation.slide_height
            slide.shapes.add_picture(image_path, left, top, width=slide_width, height=slide_height)

        # Salvar a apresentação como arquivo PPT
        presentation.save(output_ppt_path)
        print(f"Arquivo PPT salvo em: {output_ppt_path}")

    finally:
        # Remover as imagens temporárias
        print("Limpando arquivos temporários...")
        for file_name in os.listdir(temp_folder):
            file_path = os.path.join(temp_folder, file_name)
            os.remove(file_path)
        os.rmdir(temp_folder)

# Uso da função
# pdf_file_path = r"C:\Users\allan.ribeiro\Documents\funcoes_python\converter_pdf_ppt\embrapii_visão_geral.pdf"
pdf_file_path = 'Embrapii_2024_Conferência_Anual_de_Unidades_Embrapii.pdf'
output_ppt_path = "Embrapii_2024_Conferência_Anual_de_Unidades_Embrapii.pptx"
pdf_to_ppt_image(pdf_file_path, output_ppt_path)

# Requisitos:
# - PyMuPDF (pip install PyMuPDF)
# - python-pptx (pip install python-pptx)